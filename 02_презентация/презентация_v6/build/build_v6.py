# -*- coding: utf-8 -*-
"""
build_v6.py — сборка презентации «Не просто накормить» v6
Дизайн-система: FAF9F5 / F0EEE6 / 141413 / C15F3C
Шрифты: Playfair Display (заголовки/цифры), PT Sans (тело)
Сетка: 16:9, 13.333" x 7.5" (стандартный widescreen PPTX)

Структура (storyboard_v6.md):
- 18 слайдов основного показа
- 1 разделитель «Резерв»
- 12 резервных слайдов (R1-R12)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import os

# ============== ПАЛИТРА ==============
IVORY      = RGBColor(0xFA, 0xF9, 0xF5)
PANEL      = RGBColor(0xF0, 0xEE, 0xE6)
PANELWARM  = RGBColor(0xF6, 0xF1, 0xE9)
ACCENTSOFT = RGBColor(0xF5, 0xE7, 0xDE)
INK        = RGBColor(0x14, 0x14, 0x13)
BODY       = RGBColor(0x3D, 0x39, 0x29)
MUTED      = RGBColor(0x87, 0x86, 0x7F)
SOFT       = RGBColor(0x6E, 0x6A, 0x5E)
ACCENT     = RGBColor(0xC1, 0x5F, 0x3C)
ACCENTDEEP = RGBColor(0xB8, 0x55, 0x2F)
LINE       = RGBColor(0xE3, 0xDC, 0xCE)
GRAY       = RGBColor(0xB7, 0xAF, 0xA0)

SERIF = "Playfair Display"
SANS  = "PT Sans"

# ============== РАЗМЕРЫ ==============
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============== HELPERS ==============

def add_textbox(slide, x, y, w, h, text, *,
                font=SANS, size=14, bold=False, italic=False,
                color=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.3, word_wrap=True):
    """Добавить текстбокс с одним или несколькими run'ами."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0, corner_radius=None):
    """Добавить прямоугольник (опц. скруглённый)."""
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        # Установить радиус
        try:
            shape.adjustments[0] = corner_radius / 100.0
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width) if line_width else Pt(0.75)
    else:
        shape.line.fill.background()
    # Снять тень и эффекты
    sp = shape.shadow
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, weight_pt=0.75):
    """Горизонтальная или диагональная линия через connector."""
    from pptx.util import Emu
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_kicker(slide, label):
    """Верхняя строка: kicker + разделительная линия."""
    add_textbox(slide, Inches(0.72), Inches(0.5), Inches(8.0), Inches(0.3),
                label.upper(), font=SANS, size=11, bold=True, color=ACCENT)
    # Линия после kicker
    add_line(slide, Inches(0.72) + Inches(2.6), Inches(0.65), Inches(12.6), Inches(0.65),
             color=LINE, weight_pt=0.75)


def add_footer(slide, left_text, page_num, total=18, dark=False):
    """Подвал слайда: левый текст, правый — номер/всего."""
    color = MUTED if not dark else RGBColor(0x80, 0x7E, 0x76)
    add_textbox(slide, Inches(0.72), Inches(7.05), Inches(8.0), Inches(0.25),
                left_text, font=SANS, size=9, color=color)
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.1), Inches(0.25),
                f"{page_num} / {total}", font=SANS, size=9, color=color, align=PP_ALIGN.RIGHT)
    # Акцентный номер страницы
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(0.4), Inches(0.25),
                str(page_num), font=SANS, size=9, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


def set_bg(slide, color=IVORY):
    """Фон слайда."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, x, y, w=Inches(0.8), h=Inches(0.05), color=ACCENT):
    """Терракотовая черта-маркер."""
    return add_rect(slide, x, y, w, h, fill=color)


# ============== СЛАЙД 1: ТИТУЛ ==============
def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)

    # Верхняя узкая панель
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.3), fill=PANEL)

    # Логотип-эмблема (8-конечная звезда — без клипарта, геометрия)
    star_size = Inches(0.32)
    sx = Inches(0.72)
    sy = Inches(0.85)
    star = slide.shapes.add_shape(MSO_SHAPE.OVAL, sx, sy, star_size, star_size)
    star.fill.solid()
    star.fill.fore_color.rgb = ACCENT
    star.line.fill.background()
    # Крест внутри
    add_textbox(slide, sx, sy, star_size, star_size, "✦",
                font=SERIF, size=18, bold=True, color=IVORY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Kicker линия
    add_textbox(slide, Inches(1.2), Inches(0.92), Inches(11.0), Inches(0.3),
                "ПРОСТРАНСТВО НОВЫХ ИДЕЙ 2.0  ·  САНКТ-ПЕТЕРБУРГ  ·  25–27 АВГУСТА 2026",
                font=SANS, size=10, bold=True, color=SOFT)

    add_textbox(slide, Inches(0.72), Inches(1.65), Inches(12.0), Inches(0.3),
                "ДИСКУССИЯ  ·  26 АВГУСТА  ·  14:00–15:30  ·  ПЛОЩАДКА № 2 — ДСО „ТЕСОВЫЙ БЕРЕГ“",
                font=SANS, size=11, bold=True, color=ACCENT, line_spacing=1.0)

    # Главный заголовок
    add_textbox(slide, Inches(0.72), Inches(2.25), Inches(12.0), Inches(1.2),
                "Не просто накормить",
                font=SERIF, size=64, bold=True, color=INK, line_spacing=1.05)

    # Подзаголовок
    add_textbox(slide, Inches(0.72), Inches(3.55), Inches(12.0), Inches(0.5),
                "Как организовать безопасное, достойное и вариативное питание",
                font=SERIF, size=22, italic=True, color=BODY, line_spacing=1.3)
    add_textbox(slide, Inches(0.72), Inches(3.95), Inches(12.0), Inches(0.5),
                "людей с психическими нарушениями в домах социального обслуживания",
                font=SERIF, size=22, italic=True, color=BODY, line_spacing=1.3)

    # Терракотовая черта
    add_accent_bar(slide, Inches(0.72), Inches(4.9), w=Inches(0.9), h=Inches(0.05))

    # Модераторы
    add_textbox(slide, Inches(0.72), Inches(5.1), Inches(12.0), Inches(0.3),
                "МОДЕРАТОРЫ ДИСКУССИИ",
                font=SANS, size=10, bold=True, color=MUTED)

    add_textbox(slide, Inches(0.72), Inches(5.5), Inches(6.0), Inches(0.5),
                ["Евгений Чистяков", "директор СПб ГАСУСОН «ДСО „Серафимовский“»"],
                font=SANS, size=12, color=BODY, line_spacing=1.4)
    # bold на имени
    tb = slide.shapes[-1]
    p = tb.text_frame.paragraphs[0]
    p.runs[0].font.bold = True

    add_textbox(slide, Inches(7.0), Inches(5.5), Inches(6.0), Inches(0.5),
                ["Тимур Нурбаев", "директор СПб ГБСУСОН «ДСО „Тесовый берег“»"],
                font=SANS, size=12, color=BODY, line_spacing=1.4)
    tb = slide.shapes[-1]
    p = tb.text_frame.paragraphs[0]
    p.runs[0].font.bold = True

    # Footer
    add_textbox(slide, Inches(0.72), Inches(7.1), Inches(10.0), Inches(0.25),
                "Министерство труда и социальной защиты РФ  ·  Комитет по социальной политике Санкт-Петербурга",
                font=SANS, size=9, color=MUTED)
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(1.1), Inches(0.25),
                "1 / 18", font=SANS, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(0.3), Inches(0.25),
                "1", font=SANS, size=9, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


# ============== СЛАЙД 2: ПЯТЬ ИЗМЕРЕНИЙ ОДНОГО ОБЕДА ==============
def slide_02_five_dimensions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Конфликт за одним столом")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "За одним обедом — пять вещей одновременно",
                font=SERIF, size=30, bold=True, color=INK, line_spacing=1.1)

    # 5 столбцов
    labels = ["Безопасность", "Достоинство", "Предпочтение", "Возможности кухни", "Документ проверки"]
    descs = [
        "текстуры, лекарства, риск аспирации",
        "уважение к жителю, имя на тарелке",
        "свой голос — словом, жестом, биографией",
        "что реально приготовить сегодня",
        "журнал, карта, согласование"
    ]
    start_x = 0.72
    col_w = 2.4
    gap = 0.05
    for i, (lab, desc) in enumerate(zip(labels, descs)):
        x = Inches(start_x + i * (col_w + gap))
        # Панель
        add_rect(slide, x, Inches(2.05), Inches(col_w), Inches(2.9),
                 fill=PANEL, corner_radius=12)
        # Цифра
        add_textbox(slide, x, Inches(2.25), Inches(col_w), Inches(0.5),
                    f"0{i+1}", font=SERIF, size=36, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        # Метка
        add_textbox(slide, x + Inches(0.15), Inches(3.1), Inches(col_w - 0.3), Inches(0.7),
                    lab, font=SERIF, size=15, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, line_spacing=1.2)
        # Описание
        add_textbox(slide, x + Inches(0.15), Inches(4.0), Inches(col_w - 0.3), Inches(0.9),
                    desc, font=SANS, size=11, color=SOFT,
                    align=PP_ALIGN.CENTER, line_spacing=1.35)

    # Подпись снизу
    add_textbox(slide, Inches(0.72), Inches(5.3), Inches(12.0), Inches(0.5),
                "Каждое измерение — отдельный документ, отдельный участник, отдельный риск.",
                font=SERIF, size=15, italic=True, color=SOFT, align=PP_ALIGN.CENTER)

    # Терракотовая черта
    add_accent_bar(slide, Inches(6.16), Inches(6.05), w=Inches(1.0), h=Inches(0.04))

    add_footer(slide, "Доклад, разделы 4–8; сценарий 0:00–2:00", 2)


# ============== СЛАЙД 3: ДИСФАГИЯ ==============
def slide_03_dysphagia(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Почему кухня не выбирает")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Дисфагия встречается часто.",
                font=SERIF, size=30, bold=True, color=INK)

    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.4),
                "Именно поэтому текстуру назначает врач, а не повар.",
                font=SERIF, size=22, italic=True, color=SOFT)

    # Левая часть — большая цифра
    add_textbox(slide, Inches(0.72), Inches(2.7), Inches(7.0), Inches(1.8),
                "21,9–69,5 %",
                font=SERIF, size=56, bold=True, color=INK, line_spacing=1.0)

    add_textbox(slide, Inches(0.72), Inches(3.9), Inches(7.0), Inches(0.6),
                "распространённость в отдельных изученных клинических группах",
                font=SANS, size=13, color=SOFT, line_spacing=1.3)

    # Правая часть — оговорка
    add_rect(slide, Inches(7.6), Inches(2.6), Inches(5.0), Inches(2.6),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(7.85), Inches(2.85), Inches(4.6), Inches(0.4),
                "ВАЖНАЯ ОГОВОРКА",
                font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(7.85), Inches(3.25), Inches(4.6), Inches(2.0),
                ["Переносить этот диапазон",
                 "на все российские ДСО нельзя.",
                 "",
                 "Решение о текстуре —",
                 "за врачом."],
                font=SANS, size=13, color=INK, line_spacing=1.45)

    # Нижняя подпись
    add_textbox(slide, Inches(0.72), Inches(6.05), Inches(12.0), Inches(0.4),
                "Источник: обзор эпидемиологии нарушений глотания  ·  доклад, раздел 10  ·  пневмония при тяжёлых психических расстройствах — ×7 популяционной смертности",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "S013  ·  обзор эпидемиологии нарушений глотания", 3)


# ============== СЛАЙД 4: 2 190 ПРИЁМОВ В ГОД ==============
def slide_04_2190(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Масштаб повседневности")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Самые частые решения, которые дом принимает за жителя",
                font=SERIF, size=26, bold=True, color=INK)

    # Огромная цифра (неразрывный пробел чтобы "2 190" не разрывалось)
    add_textbox(slide, Inches(0.72), Inches(2.1), Inches(8.0), Inches(2.8),
                "2\u00a0190",
                font=SERIF, size=110, bold=True, color=INK, line_spacing=1.0)

    # Под цифрой
    add_textbox(slide, Inches(0.72), Inches(4.8), Inches(7.5), Inches(0.5),
                "приёмов пищи в год за одного человека",
                font=SANS, size=18, bold=True, color=SOFT)

    add_textbox(slide, Inches(0.72), Inches(5.25), Inches(7.5), Inches(0.4),
                "6 приёмов × 365 дней  =  2\u00a0190",
                font=SANS, size=14, color=MUTED)

    # Правая колонка: 6 приёмов
    meals = [
        ("завтрак", 0.05),
        ("2-й завтрак", 0.27),
        ("обед", 0.49),
        ("полдник", 0.05),
        ("ужин", 0.27),
        ("вечерний", 0.49),
    ]
    for i, (meal, _) in enumerate(meals):
        y = Inches(2.1 + i * 0.62)
        # Кружок
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), y, Inches(0.3), Inches(0.3))
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
        c.line.fill.background()
        # Текст
        add_textbox(slide, Inches(9.0), y + Inches(0.02), Inches(3.5), Inches(0.3),
                    meal, font=SERIF, size=18, bold=True, color=INK)

    add_textbox(slide, Inches(8.5), Inches(5.95), Inches(4.5), Inches(0.4),
                "шесть раз в день  ·  365 дней  ·  одна жизнь",
                font=SERIF, size=14, italic=True, color=SOFT)

    add_footer(slide, "S030  ·  Комитет по социальной политике СПб, 22.07.2025", 4)


# ============== СЛАЙД 5: 74 % ==============
def slide_05_74(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Контекст · недееспособность")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Порядка трёх из четырёх — недееспособны.",
                font=SERIF, size=30, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.55), Inches(12.0), Inches(0.5),
                "Их мнение по ст. 29 ГК обязаны выяснять и учитывать.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # Donut
    cx, cy, r = Inches(3.0), Inches(3.7), Inches(1.5)
    donut_w = Inches(0.45)
    # Серый круг
    gray = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    gray.fill.solid()
    gray.fill.fore_color.rgb = PANEL
    gray.line.fill.background()
    # Терракотовый "arc" — реализован как дуговая полоса через OVAL с заполнением
    accent = slide.shapes.add_shape(MSO_SHAPE.PIE, cx - r, cy - r, r * 2, r * 2)
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    # Регулировка старта/угла дуги
    try:
        accent.adjustments[0] = 0  # start angle
        accent.adjustments[1] = 266  # end angle 74% of 360 = 266
    except Exception:
        pass
    # Внутренний круг (дырка)
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r + donut_w, cy - r + donut_w,
                                   (r - donut_w) * 2, (r - donut_w) * 2)
    inner.fill.solid()
    inner.fill.fore_color.rgb = IVORY
    inner.line.fill.background()
    # Текст в центре
    add_textbox(slide, cx - r, cy - Inches(0.4), r * 2, Inches(0.7),
                "74 %", font=SERIF, size=46, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, cx - r, cy + Inches(0.3), r * 2, Inches(0.4),
                "недееспособны", font=SANS, size=12, color=MUTED,
                align=PP_ALIGN.CENTER)

    # Правая часть — пояснение
    add_textbox(slide, Inches(5.5), Inches(2.4), Inches(7.2), Inches(0.4),
                "ПОВСЕДНЕВНОЕ ПРЕДПОЧТЕНИЕ — НЕ СДЕЛКА",
                font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(5.5), Inches(2.75), Inches(7.2), Inches(2.5),
                ["Сделки совершает опекун (ст. 29 ГК).",
                 "Но повседневное предпочтение — это мнение,",
                 "а мнение обязаны выяснять.",
                 "",
                 "Кто не говорит — показ, наблюдение, биография."],
                font=SANS, size=14, color=INK, line_spacing=1.5)

    # Оговорка
    add_textbox(slide, Inches(0.72), Inches(6.3), Inches(12.0), Inches(0.5),
                ["Расчёт проекта «Если быть точным», данные за 2024 г.  ·  не статотчётность Минтруда  ·  в вашем доме цифра будет другой."],
                font=SANS, size=11, italic=True, color=SOFT, line_spacing=1.4)

    add_footer(slide, "S004  ·  проект «Если быть точным», публикация 09.04.2026  ·  ст. 29 ГК, 48-ФЗ", 5)


# ============== СЛАЙД 6: НОЧНОЙ ИНТЕРВАЛ ==============
def slide_06_night(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Физиология · ужин → завтрак")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Ужин 17:30  →  завтрак 8:30  =  15 часов без еды",
                font=SERIF, size=26, bold=True, color=INK)

    add_textbox(slide, Inches(0.72), Inches(1.55), Inches(12.0), Inches(0.4),
                "Цель нашего дома — не больше 13 часов.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # Лента суток
    lane_x_in = 0.72
    lane_y_in = 3.2
    lane_w_in = 11.9
    lane_h_in = 0.7
    lane_x = Inches(lane_x_in)
    lane_y = Inches(lane_y_in)
    lane_w = Inches(lane_w_in)
    lane_h = Inches(lane_h_in)

    # Заштрихованный "день" (6:00–22:00)
    day_start_pct = 6/24
    day_end_pct = 22/24
    # Фон
    add_rect(slide, lane_x, lane_y, lane_w, lane_h, fill=PANEL, corner_radius=8)
    # Дневная часть (светлая)
    day_w = Inches(lane_w_in * (day_end_pct - day_start_pct))
    day_x = Inches(lane_x_in + lane_w_in * day_start_pct)
    add_rect(slide, day_x, lane_y, day_w, lane_h, fill=IVORY, line_color=LINE, line_width=0.5)

    # Метки времени
    for h, lbl in [(0, "00"), (6, "06"), (12, "12"), (18, "18"), (24, "24")]:
        x = Inches(lane_x_in + lane_w_in * h / 24 - 0.2)
        add_textbox(slide, x, lane_y + lane_h + Inches(0.05), Inches(0.4), Inches(0.3),
                    lbl, font=SANS, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # Точки приёмов (завтрак 8:30, обед 13:00, ужин 17:30)
    # Реальные позиции по ленте
    meal_pcts = [0.354, 0.542, 0.729]  # 8:30, 13:00, 17:30
    meal_xs_in = [lane_x_in + lane_w_in * p for p in meal_pcts]
    for mx_in in meal_xs_in:
        x = Inches(mx_in - 0.09)
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, lane_y + Inches(0.22), Inches(0.18), Inches(0.26))
        c.fill.solid()
        c.fill.fore_color.rgb = INK
        c.line.color.rgb = IVORY
        c.line.width = Pt(2)

    # Подписи приёмов — СТРОГО под точками
    meal_labels = [("завтрак", 0.354), ("обед", 0.542), ("ужин", 0.729)]
    for name, pct in meal_labels:
        cx = lane_x_in + lane_w_in * pct
        # Текст шириной 1.4", по центру под точкой
        add_textbox(slide, Inches(cx - 0.7), lane_y + lane_h + Inches(0.4),
                    Inches(1.4), Inches(0.3),
                    name, font=SANS, size=11, color=INK, align=PP_ALIGN.CENTER)
        # Время под подписью (мелко)
        time_label = {"завтрак": "8:30", "обед": "13:00", "ужин": "17:30"}[name]
        add_textbox(slide, Inches(cx - 0.7), lane_y + lane_h + Inches(0.7),
                    Inches(1.4), Inches(0.3),
                    time_label, font=SANS, size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # Двойная дуга-стрелка «15 ч через полночь» НАД лентой
    # От точки ужина (17:30, 0.729) к точке завтрака (8:30 следующего дня, 1.354)
    # Визуально: арка поднимается над лентой, с тонкой линией и стрелками на концах
    dinner_x = lane_x_in + lane_w_in * 0.729
    breakfast_x = lane_x_in + lane_w_in * 0.354  # 8:30 следующего дня → рисуем как «возврат к началу»
    # Рисуем арку: от (dinner_x, lane_y - 0.3) до (breakfast_x, lane_y - 0.3) с подъёмом
    # Но breakfast_x = 0.72 + 11.9*0.354 = 4.93, а dinner_x = 0.72 + 11.9*0.729 = 9.39
    # Значит арка идёт СПРАВА НАЛЕВО — что неестественно для стрелки
    # Решение: нарисовать арку с разрывом — два сегмента:
    #   1) от ужина (17:30) до 24:00 — обычный горизонтальный сегмент
    #   2) от 00:00 до 8:30 — обычный горизонтальный сегмент
    # С подписью "через полночь" посередине

    arc_y = lane_y - Inches(0.55)
    arc_h = Inches(0.45)
    # Сегмент 1: 17:30 → 24:00
    seg1_x1 = dinner_x
    seg1_x2 = lane_x_in + lane_w_in
    # Сегмент 2: 00:00 → 8:30
    seg2_x1 = lane_x_in
    seg2_x2 = breakfast_x

    # Линии-сегменты
    line1 = slide.shapes.add_connector(1, Inches(seg1_x1), arc_y, Inches(seg1_x2), arc_y)
    line1.line.color.rgb = ACCENT
    line1.line.width = Pt(2.5)
    line2 = slide.shapes.add_connector(1, Inches(seg2_x1), arc_y, Inches(seg2_x2), arc_y)
    line2.line.color.rgb = ACCENT
    line2.line.width = Pt(2.5)

    # Стрелки на концах (в точках ужина и завтрака)
    # Стрелка в начале (завтрак)
    a1 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(seg2_x2) - Inches(0.05),
                                  arc_y - Inches(0.07), Inches(0.18), Inches(0.14))
    a1.fill.solid(); a1.fill.fore_color.rgb = ACCENT
    a1.line.fill.background()
    # Стрелка в конце (ужин)
    a2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(seg1_x1) - Inches(0.13),
                                  arc_y - Inches(0.07), Inches(0.18), Inches(0.14))
    a2.fill.solid(); a2.fill.fore_color.rgb = ACCENT
    a2.line.fill.background()

    # Метка «15 ч» посередине (между двумя сегментами)
    mid_x = (seg1_x2 + seg2_x1) / 2
    add_textbox(slide, Inches(mid_x - 0.9), arc_y - Inches(0.55), Inches(1.8), Inches(0.35),
                "через полночь  ·  15 ч",
                font=SERIF, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    # Тонкая вертикаль-разделитель (полночь) — пунктирная
    midnight_x = lane_x_in + lane_w_in
    add_line(slide, Inches(midnight_x), arc_y, Inches(midnight_x), arc_y + Inches(0.5),
             color=ACCENT, weight_pt=1.0)
    add_textbox(slide, Inches(midnight_x - 0.5), arc_y - Inches(0.22), Inches(1.0), Inches(0.25),
                "полночь", font=SANS, size=8, color=MUTED, align=PP_ALIGN.CENTER)

    # Сравнение
    add_rect(slide, Inches(0.72), Inches(4.8), Inches(5.8), Inches(1.8),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(5.0), Inches(5.5), Inches(0.4),
                "ОРИЕНТИРЫ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(5.4), Inches(5.5), Inches(1.0),
                ["≤ 11 ч  —  Швеция",
                 "≤ 13 ч  —  цель нашего дома",
                 "≤ 14 ч  —  зарубежные ориентиры (диапазон)"],
                font=SANS, size=14, color=INK, line_spacing=1.5)

    # Оговорка
    add_rect(slide, Inches(6.85), Inches(4.8), Inches(5.8), Inches(1.8),
             fill=ACCENTSOFT, corner_radius=14)
    add_textbox(slide, Inches(7.1), Inches(5.0), Inches(5.5), Inches(0.4),
                "ОГОВОРКА", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(7.1), Inches(5.4), Inches(5.5), Inches(1.2),
                ["15 ч — арифметика одного расписания,",
                 "не статистика по российским ДСО.",
                 "Швеция/США — другие системы регулирования."],
                font=SANS, size=12, color=INK, line_spacing=1.45)

    add_footer(slide, "Доклад, раздел 7.4, 26  ·  Швеция ≤11 ч  ·  сравнение ограничено", 6)


# ============== СЛАЙД 7: КОМПЛЕКТ, НЕ ЛОЗУНГ ==============
def slide_07_folder(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Право · не лозунг, а папка")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Проверяющему нужна не фраза «запрета нет».",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "А положение, журнал, карты, согласование с учредителем.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # Центральная "папка" — панель с пятью вкладками
    add_rect(slide, Inches(0.72), Inches(2.4), Inches(11.9), Inches(3.0),
             fill=PANEL, corner_radius=14)

    add_textbox(slide, Inches(1.0), Inches(2.55), Inches(11.5), Inches(0.4),
                "МИНИМАЛЬНЫЙ КОМПЛЕКТ ДОКУМЕНТОВ",
                font=SANS, size=11, bold=True, color=ACCENTDEEP)

    items = [
        ("①", "Положение", "о питании, в редакции с учётом вариативности"),
        ("②", "Меню", "с вариантами и парами эквивалентов"),
        ("③", "Журнал заказа", "что заказал житель, кто принял"),
        ("④", "Журнал замен", "когда варианта не оказалось, что выдали"),
        ("⑤", "Пищевые карты", "профиль, текстура, исключения, предпочтения"),
        ("⑥", "Согласование", "с учредителем — для сведения или по регламенту"),
    ]

    grid_x = [0.95, 5.05, 9.15]
    grid_y = [3.05, 4.5]
    for i, (num, title, desc) in enumerate(items):
        col = i % 3
        row = i // 3
        x = Inches(grid_x[col])
        y = Inches(grid_y[row])
        add_textbox(slide, x, y, Inches(0.4), Inches(0.5),
                    num, font=SERIF, size=24, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(0.45), y, Inches(3.6), Inches(0.4),
                    title, font=SERIF, size=15, bold=True, color=INK)
        add_textbox(slide, x + Inches(0.45), y + Inches(0.4), Inches(3.6), Inches(0.7),
                    desc, font=SANS, size=11, color=SOFT, line_spacing=1.35)

    # Нижняя полоса
    add_accent_bar(slide, Inches(0.72), Inches(5.65), w=Inches(0.9), h=Inches(0.04))
    add_textbox(slide, Inches(0.72), Inches(5.85), Inches(12.0), Inches(0.5),
                "Федеральное право не запрещает и не обязывает — это нужно закрепить локально.",
                font=SERIF, size=16, italic=True, color=SOFT, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "Источники: 442-ФЗ ст. 9, 16, 32 ч. 4  ·  3185-1 ст. 5, 37, 43  ·  ГК ст. 29  ·  доклад, разделы 31, 32",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Без комплекта та же практика выглядит как отступление от раскладки", 7)


# ============== СЛАЙД 8: ИРКУТСКАЯ ОБЛАСТЬ ==============
def slide_08_irkutsk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Российская практика")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Иркутская область масштабирует выбор на регион.",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Не один эксперимент — контур для учредителя.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # Левая карточка — метка
    add_rect(slide, Inches(0.72), Inches(2.6), Inches(4.5), Inches(3.4),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.72), Inches(2.85), Inches(4.5), Inches(0.5),
                "ИРКУТСКАЯ ОБЛАСТЬ", font=SANS, size=11, bold=True, color=ACCENTDEEP,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.72), Inches(3.4), Inches(4.5), Inches(1.5),
                "вариативное\nпитание\nв стационарных\nучреждениях",
                font=SERIF, size=24, bold=True, color=INK,
                align=PP_ALIGN.CENTER, line_spacing=1.15)

    add_textbox(slide, Inches(0.72), Inches(5.4), Inches(4.5), Inches(0.4),
                "публикация 14.02.2025", font=SANS, size=11, color=SOFT,
                align=PP_ALIGN.CENTER)

    # Правая часть — содержание
    add_textbox(slide, Inches(5.6), Inches(2.6), Inches(7.0), Inches(0.4),
                "ЧТО СДЕЛАНО", font=SANS, size=11, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(5.6), Inches(2.95), Inches(7.0), Inches(1.5),
                ["Выбор по основным приёмам пищи:",
                 "  ·  два первых блюда",
                 "  ·  два вторых блюда",
                 "  ·  два гарнира"],
                font=SERIF, size=18, color=INK, line_spacing=1.4)

    add_textbox(slide, Inches(5.6), Inches(4.7), Inches(7.0), Inches(0.4),
                "ЧТО ЭТО ЗНАЧИТ ДЛЯ ДИРЕКТОРА", font=SANS, size=11, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(5.6), Inches(5.05), Inches(7.0), Inches(1.2),
                ["Не «дайте денег — сделаем».",
                 "Контур, который можно показать учредителю",
                 "при разговоре о масштабировании."],
                font=SANS, size=13, color=SOFT, line_spacing=1.45)

    # Оговорка
    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "Источник S018  ·  масштаб области как аудит не подтверждён; показано как контур, а не как результат",
                font=SANS, size=10, italic=True, color=MUTED)

    add_footer(slide, "Успенский, Болотнинский, Усть-Илимск, Серафимовский — в резерве R3–R6", 8)


# ============== СЛАЙД 9: СЕРАФИМОВСКИЙ — ЗАКАЗ НАКАНУНЕ ==============
def slide_09_serafimovsky(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Кейс · наш дом")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Житель заказывает ужин накануне.",
                font=SERIF, size=30, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "К вечеру известно, что готовить.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # "Лента дня" — горизонтальная с двумя точками
    add_rect(slide, Inches(0.72), Inches(2.7), Inches(11.9), Inches(1.2),
             fill=PANEL, corner_radius=14)

    add_textbox(slide, Inches(0.72), Inches(2.85), Inches(11.9), Inches(0.4),
                "ВЧЕРА  ·  СЕГОДНЯ", font=SANS, size=10, bold=True, color=ACCENTDEEP,
                align=PP_ALIGN.CENTER)

    # Два "окна"
    add_textbox(slide, Inches(1.5), Inches(3.3), Inches(4.5), Inches(0.5),
                "18:00", font=SERIF, size=30, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(3.8), Inches(4.5), Inches(0.4),
                "принимаем заказ на ужин", font=SANS, size=13, color=SOFT,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(7.3), Inches(3.3), Inches(4.5), Inches(0.5),
                "08:00", font=SERIF, size=30, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.3), Inches(3.8), Inches(4.5), Inches(0.4),
                "подача на завтрак", font=SANS, size=13, color=SOFT,
                align=PP_ALIGN.CENTER)

    # Стрелка
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), Inches(3.35),
                                    Inches(1.2), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()

    # Нижняя часть
    add_textbox(slide, Inches(0.72), Inches(4.4), Inches(12.0), Inches(0.4),
                "ЧТО ЭТО ДАЛО", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.72), Inches(4.75), Inches(12.0), Inches(1.0),
                ["Жители приняли выбор с энтузиазмом. Директор — я — рассказываю и про цену, и про пределы."],
                font=SERIF, size=18, color=INK, line_spacing=1.4)

    add_textbox(slide, Inches(0.72), Inches(5.9), Inches(12.0), Inches(0.4),
                "Старт не был бесплатным  —  дополнительное оборудование и ставки пищеблока.",
                font=SANS, size=12, bold=True, italic=True, color=SOFT)

    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "Источник S030  ·  gov.spb.ru, 22.07.2025  ·  СПб ГАСУСОН «ДСО „Серафимовский“»",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Кейс — часть практики, не доказательство эффекта", 9)


# ============== СЛАЙД 10: ПРИКАЗ № 124 ==============
def slide_10_order(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Локальный норматив")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Приказ № 124  от  10.03.2026",
                font=SERIF, size=44, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.7), Inches(12.0), Inches(0.5),
                "Заказное питание вошло в правила внутреннего распорядка жителей ДСО „Серафимовский“.",
                font=SERIF, size=18, italic=True, color=SOFT, line_spacing=1.3)

    # Левая карточка
    add_rect(slide, Inches(0.72), Inches(2.6), Inches(6.0), Inches(3.4),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.8), Inches(5.6), Inches(0.4),
                "ЧТО ИМЕННО ЗАКРЕПЛЕНО", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(3.2), Inches(5.6), Inches(2.7),
                ["Заказное питание (вариативное меню)",
                 "как часть правил внутреннего распорядка,",
                 "а не «эксперимент по новости».",
                 "",
                 "Включая маломобильных:",
                 "  ·  приём пищи в жилых комнатах",
                 "  ·  помощь в кормлении"],
                font=SANS, size=14, color=INK, line_spacing=1.5)

    # Правая карточка
    add_rect(slide, Inches(7.0), Inches(2.6), Inches(5.6), Inches(3.4),
             fill=ACCENTSOFT, corner_radius=14)
    add_textbox(slide, Inches(7.25), Inches(2.8), Inches(5.2), Inches(0.4),
                "ПОЧЕМУ ЭТО ОПОРА", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(7.25), Inches(3.2), Inches(5.2), Inches(2.7),
                ["«Это правила дома, утверждённые приказом» —",
                 "аргумент, который пережил смену настроения.",
                 "",
                 "Проверяющему — документ, не слова.",
                 "Учредителю — порядок, не самочинность."],
                font=SANS, size=14, color=INK, line_spacing=1.5)

    # Оговорка
    add_textbox(slide, Inches(0.72), Inches(6.3), Inches(12.0), Inches(0.5),
                "Локальный акт одного учреждения  ·  не общероссийский норматив.",
                font=SANS, size=11, italic=True, bold=True, color=ACCENTDEEP)

    add_textbox(slide, Inches(0.72), Inches(6.7), Inches(12.0), Inches(0.3),
                "Источник S035  ·  www.pni9.ru  ·  правила внутреннего распорядка жителей",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Доклад, раздел 8, 17.12", 10)


# ============== СЛАЙД 11: 223-ФЗ / 44-ФЗ ==============
def slide_11_procurement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Закупки")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Один учредитель — две разные закупки.",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Вариативность не от номера закона.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # Две карточки
    add_rect(slide, Inches(0.72), Inches(2.5), Inches(5.8), Inches(3.5),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.7), Inches(5.5), Inches(0.4),
                "СЕРАФИМОВСКИЙ", font=SANS, size=11, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(3.1), Inches(5.5), Inches(0.7),
                "223-ФЗ", font=SERIF, size=48, bold=True, color=INK)
    add_textbox(slide, Inches(0.95), Inches(4.1), Inches(5.5), Inches(0.4),
                "автономное учреждение  ·  конкурсная основа",
                font=SANS, size=12, color=SOFT, line_spacing=1.4)
    add_textbox(slide, Inches(0.95), Inches(4.65), Inches(5.5), Inches(1.3),
                ["Тендеры: 32616151967, 95375953, 32616197209",
                 "tenmon.ru, tenderguru.ru, 2026"],
                font=SANS, size=11, color=MUTED, line_spacing=1.5)

    add_rect(slide, Inches(6.8), Inches(2.5), Inches(5.8), Inches(3.5),
             fill=ACCENTSOFT, corner_radius=14)
    add_textbox(slide, Inches(7.05), Inches(2.7), Inches(5.5), Inches(0.4),
                "ТЕСОВЫЙ БЕРЕГ", font=SANS, size=11, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(7.05), Inches(3.1), Inches(5.5), Inches(0.7),
                "44-ФЗ", font=SERIF, size=48, bold=True, color=INK)
    add_textbox(slide, Inches(7.05), Inches(4.1), Inches(5.5), Inches(0.4),
                "бюджетное учреждение  ·  государственная закупка",
                font=SANS, size=12, color=SOFT, line_spacing=1.4)
    add_textbox(slide, Inches(7.05), Inches(4.65), Inches(5.5), Inches(1.3),
                ["ИНН 7827661874  ·  ЕГРЮЛ",
                 "Санкт-Петербург, п. Молодёжное"],
                font=SANS, size=11, color=MUTED, line_spacing=1.5)

    # Вывод
    add_accent_bar(slide, Inches(0.72), Inches(6.3), w=Inches(0.9), h=Inches(0.04))
    add_textbox(slide, Inches(0.72), Inches(6.5), Inches(12.0), Inches(0.5),
                "Вариативность зависит не от номера закона — а от корректной организации требований, документов и исполнения.",
                font=SERIF, size=15, italic=True, color=SOFT, align=PP_ALIGN.CENTER, line_spacing=1.3)

    add_footer(slide, "Доклад, раздел 25, 29  ·  программа организаторов", 11)


# ============== СЛАЙД 12: СНАЧАЛА ВРАЧ, ПОТОМ ВЫБОР ==============
def slide_12_doctor_first(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Сначала рамка")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Сначала медицинская рамка. Потом выбор.",
                font=SERIF, size=30, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Кухня исполняет назначение, а не ставит его.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # Две горизонтальные карточки
    add_rect(slide, Inches(0.72), Inches(2.4), Inches(11.9), Inches(1.5),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.55), Inches(11.5), Inches(0.4),
                "РЕШАЕТ ВРАЧ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(2.95), Inches(11.5), Inches(0.9),
                ["Текстура  ·  стол  ·  противопоказания  ·  лекарственно-пищевые исключения",
                 "Приказы 330н, 395н;  СанПиН 2.3/2.4.4282-26, п. 56(11)"],
                font=SERIF, size=16, color=INK, line_spacing=1.4)

    add_rect(slide, Inches(0.72), Inches(4.05), Inches(11.9), Inches(1.5),
             fill=ACCENTSOFT, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(4.2), Inches(11.5), Inches(0.4),
                "РЕШАЕТ УЧРЕЖДЕНИЕ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(4.6), Inches(11.5), Inches(0.9),
                ["Сколько вариантов  ·  где выбор фиксируется  ·  что делает смена, когда вариант кончился",
                 "Приказ, положение, журнал, линия раздачи"],
                font=SERIF, size=16, color=INK, line_spacing=1.4)

    # Нижний вывод
    add_accent_bar(slide, Inches(6.16), Inches(5.85), w=Inches(1.0), h=Inches(0.04))
    add_textbox(slide, Inches(0.72), Inches(6.05), Inches(12.0), Inches(0.5),
                "Если это звучит как медицинская реформа руками повара — вы правильно закрываете папку.",
                font=SERIF, size=16, italic=True, color=SOFT, align=PP_ALIGN.CENTER, line_spacing=1.3)

    add_footer(slide, "Доклад, раздел 11, 17  ·  правило «трёх отказов» — врач сегодня", 12)


# ============== СЛАЙД 13: МОДЕЛЬ РЕШЕНИЯ ==============
def slide_13_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Модель решения")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Врач задаёт рамку. Внутри рамки — пары эквивалентов.",
                font=SERIF, size=26, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Человек выбирает. Выбор фиксируется.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # 4 блока со стрелками
    block_w = 2.6
    block_h = 2.5
    block_y = Inches(2.5)
    gap_x = 0.25
    arrow_w = 0.45
    total_w = 4 * block_w + 3 * (gap_x + arrow_w)
    start_x = (13.333 - total_w) / 2

    blocks = [
        ("①", "Мед. ограничения", "текстура, стол,\nпротивопоказания", INK),
        ("②", "Допустимые пары", "гречка/рис,\nкурица/индейка", ACCENT),
        ("③", "Выбор человека", "показ, слово,\nнаблюдение", INK),
        ("④", "Фиксация", "в журнале заказа\nс подписью смены", INK),
    ]

    for i, (num, title, body, color) in enumerate(blocks):
        x = Inches(start_x + i * (block_w + gap_x + arrow_w))
        # Карточка
        add_rect(slide, x, block_y, Inches(block_w), Inches(block_h),
                 fill=PANEL, corner_radius=12)
        # Номер
        add_textbox(slide, x, block_y + Inches(0.15), Inches(block_w), Inches(0.5),
                    num, font=SERIF, size=28, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        # Заголовок
        add_textbox(slide, x + Inches(0.1), block_y + Inches(0.75), Inches(block_w - 0.2), Inches(0.6),
                    title, font=SERIF, size=15, bold=True, color=color,
                    align=PP_ALIGN.CENTER, line_spacing=1.2)
        # Тело
        add_textbox(slide, x + Inches(0.1), block_y + Inches(1.4), Inches(block_w - 0.2), Inches(1.0),
                    body, font=SANS, size=12, color=SOFT,
                    align=PP_ALIGN.CENTER, line_spacing=1.4)

        # Стрелка после блока (кроме последнего)
        if i < 3:
            arr_x = x + Inches(block_w + 0.02)
            arr_y = block_y + Inches(1.1)
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arr_x, arr_y,
                                          Inches(arrow_w - 0.05), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT
            arr.line.fill.background()

    # Подпись под схемой
    add_rect(slide, Inches(0.72), Inches(5.4), Inches(11.9), Inches(0.8),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.72), Inches(5.55), Inches(11.9), Inches(0.4),
                "СТАНДАРТНЫЙ ВАРИАНТ ГАРАНТИРОВАН ВСЕГДА",
                font=SANS, size=11, bold=True, color=ACCENTDEEP, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.72), Inches(5.85), Inches(11.9), Inches(0.4),
                "Человек вправе выбрать стандартное блюдо или не выбирать вовсе.",
                font=SERIF, size=15, italic=True, color=SOFT, align=PP_ALIGN.CENTER)

    add_footer(slide, "Доклад, раздел 19 (пары эквивалентов), 21 (журнал)", 13)


# ============== СЛАЙД 14: ЧЕТЫРЕ КАНАЛА ВОЛИ ==============
def slide_14_channels(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Выбор не требует дееспособности")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Четыре канала выражения воли.",
                font=SERIF, size=30, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Человек, который не говорит, всё равно может выразить предпочтение.",
                font=SERIF, size=20, italic=True, color=SOFT)

    channels = [
        ("①", "Показ", "две порции перед глазами, рука сама пойдёт"),
        ("②", "Наблюдение", "что съедено, что оставлено, «влажный» голос"),
        ("③", "Биография", "что ел всю жизнь, к чему тянется"),
        ("④", "Сведения", "от тех, кто хорошо знает человека"),
    ]
    start_x = 0.72
    box_w = 2.85
    gap = 0.1
    for i, (num, title, body) in enumerate(channels):
        x = Inches(start_x + i * (box_w + gap))
        add_rect(slide, x, Inches(2.4), Inches(box_w), Inches(3.0),
                 fill=PANEL, corner_radius=14)
        add_textbox(slide, x, Inches(2.6), Inches(box_w), Inches(0.6),
                    num, font=SERIF, size=36, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(3.4), Inches(box_w - 0.3), Inches(0.5),
                    title, font=SERIF, size=20, bold=True, color=INK,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(3.95), Inches(box_w - 0.3), Inches(1.3),
                    body, font=SANS, size=12, color=SOFT,
                    align=PP_ALIGN.CENTER, line_spacing=1.4)

    # Оговорка
    add_rect(slide, Inches(0.72), Inches(5.7), Inches(11.9), Inches(0.8),
             fill=ACCENTSOFT, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(5.85), Inches(11.5), Inches(0.4),
                "ОГОВОРКА", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(6.15), Inches(11.5), Inches(0.4),
                "Наблюдение — это сигнал, а не автоматическое согласие. Решение всё равно документируется.",
                font=SANS, size=12, color=INK, line_spacing=1.4)

    add_footer(slide, "Ст. 29 ГК  ·  Замечание общего порядка № 1 к ст. 12 КПИ", 14)


# ============== СЛАЙД 15: ПЯТЬ ШАГОВ ДИРЕКТОРА ==============
def slide_15_five_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Что сделать в понедельник")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Пилот на 90 дней. Одно отделение.",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.4),
                "Один обед  ·  одно отделение  ·  одна неделя",
                font=SERIF, size=18, italic=True, color=ACCENTDEEP)

    # 5 шагов
    steps = [
        ("①", "Замер", "7 дней\nноль рублей"),
        ("②", "Приказ", "одна строка\nоб ответственном"),
        ("③", "Обучение смены", "15 минут\nпоказать журнал"),
        ("④", "Пилот", "90 дней\nодно отделение"),
        ("⑤", "Решение о масштабе", "с письмом\nучредителю"),
    ]
    start_x = 0.72
    col_w = 2.4
    gap = 0.08
    for i, (num, title, body) in enumerate(steps):
        x = Inches(start_x + i * (col_w + gap))
        fill = ACCENT if i == 3 else PANEL
        add_rect(slide, x, Inches(2.5), Inches(col_w), Inches(2.7),
                 fill=fill, corner_radius=14)
        col_num = IVORY if i == 3 else ACCENT
        col_title = IVORY if i == 3 else INK
        col_body = IVORY if i == 3 else SOFT
        add_textbox(slide, x, Inches(2.7), Inches(col_w), Inches(0.6),
                    num, font=SERIF, size=32, bold=True, color=col_num,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(3.45), Inches(col_w - 0.3), Inches(0.7),
                    title, font=SERIF, size=15, bold=True, color=col_title,
                    align=PP_ALIGN.CENTER, line_spacing=1.2)
        add_textbox(slide, x + Inches(0.15), Inches(4.2), Inches(col_w - 0.3), Inches(0.9),
                    body, font=SANS, size=11, color=col_body,
                    align=PP_ALIGN.CENTER, line_spacing=1.4)

    # Пятиугольник измерений (низ)
    add_rect(slide, Inches(0.72), Inches(5.5), Inches(11.9), Inches(1.2),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(5.65), Inches(11.5), Inches(0.4),
                "ПЯТЬ ИЗМЕРЯЕМЫХ ПОКАЗАТЕЛЕЙ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(5.95), Inches(11.5), Inches(0.7),
                ["① пищевая карта заполнена   ② доля приёмов с реальным выбором   ③ исполненные заказы   ④ отказы/замены/отходы   ⑤ события риска"],
                font=SANS, size=11, color=INK, line_spacing=1.4)
    add_textbox(slide, Inches(0.95), Inches(6.35), Inches(11.5), Inches(0.3),
                "Процент выбравших — не KPI.",
                font=SERIF, size=12, bold=True, italic=True, color=ACCENTDEEP, align=PP_ALIGN.LEFT)

    add_footer(slide, "Доклад, раздел 22 (панель чисел), 35 (пилот)", 15)


# ============== СЛАЙД 16: ТРИ ОШИБКИ ==============
def slide_16_three_errors(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Красные линии")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Три ошибки дороже бездействия.",
                font=SERIF, size=30, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Все три лечатся процедурой.",
                font=SERIF, size=20, italic=True, color=SOFT)

    # 3 карточки
    errors = [
        ("①", "Имитация выбора",
         "«Выбирайте!» — при одном блюде на линии.",
         "Журнал врёт и проверяющему, и жителю.",
         "Переписать журнал. Сократить позиции до реальных."),
        ("②", "Выбор без гарантии",
         "Житель выбрал Б — Б закончилось.",
         "Нет стандартного — есть лотерея.",
         "Стандартный вариант всегда в наличии."),
        ("③", "«Лечебность» руками кухни",
         "Повар назначает диету.",
         "Диагноз повара — не диагноз.",
         "Назначение делает врач, кухня исполняет."),
    ]
    start_x = 0.72
    col_w = 3.95
    gap = 0.1
    for i, (num, title, what, why, fix) in enumerate(errors):
        x = Inches(start_x + i * (col_w + gap))
        add_rect(slide, x, Inches(2.4), Inches(col_w), Inches(4.0),
                 fill=PANELWARM, line_color=LINE, line_width=0.5, corner_radius=14)
        # Номер
        add_textbox(slide, x + Inches(0.3), Inches(2.6), Inches(0.5), Inches(0.5),
                    num, font=SERIF, size=24, bold=True, color=ACCENTDEEP)
        # Заголовок
        add_textbox(slide, x + Inches(0.3), Inches(3.1), Inches(col_w - 0.6), Inches(0.8),
                    title, font=SERIF, size=15, bold=True, color=INK, line_spacing=1.2)
        # Что делают
        add_textbox(slide, x + Inches(0.3), Inches(4.0), Inches(col_w - 0.6), Inches(0.8),
                    what, font=SANS, size=12, color=BODY, line_spacing=1.4)
        # Почему плохо
        add_textbox(slide, x + Inches(0.3), Inches(4.85), Inches(col_w - 0.6), Inches(0.7),
                    why, font=SANS, size=11, italic=True, color=ACCENTDEEP, line_spacing=1.4)
        # Как лечится
        add_rect(slide, x + Inches(0.3), Inches(5.65), Inches(col_w - 0.6), Inches(0.6),
                 fill=IVORY, corner_radius=8)
        add_textbox(slide, x + Inches(0.45), Inches(5.7), Inches(col_w - 0.9), Inches(0.5),
                    fix, font=SANS, size=11, bold=True, color=INK, line_spacing=1.3)

    add_footer(slide, "Доклад, раздел 17, 22  ·  сценарий 12:30–16:30", 16)


# ============== СЛАЙД 17: 01.09.2026 ==============
def slide_17_sanpin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Новый санитарный контекст")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Через шесть дней после нашей встречи",
                font=SERIF, size=22, italic=True, color=SOFT)

    # Огромная дата слева
    add_textbox(slide, Inches(0.72), Inches(1.65), Inches(6.8), Inches(2.0),
                "01.09.2026",
                font=SERIF, size=80, bold=True, color=INK, line_spacing=1.0)

    # Под числом
    add_textbox(slide, Inches(0.72), Inches(3.55), Inches(6.8), Inches(0.5),
                "вступает СанПиН 2.3/2.4.4282-26",
                font=SERIF, size=18, italic=True, color=ACCENTDEEP)

    # Правая часть — три строки
    add_textbox(slide, Inches(7.5), Inches(1.65), Inches(5.2), Inches(0.4),
                "ЧТО ЭТО ЗНАЧИТ ДЛЯ ВАШЕГО ДОМА", font=SANS, size=10, bold=True, color=ACCENTDEEP)

    rows = [
        ("Пункт 56(11):",
         "не менее 3-х приёмов, в т. ч. диетическое по показаниям"),
        ("Число вариантов блюда:",
         "ни в одной редакции правил не ограничено"),
        ("Журналы августа:",
         "законны как заполнены — задним числом не переписывать"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.1 + i * 0.9)
        add_textbox(slide, Inches(7.5), y, Inches(5.2), Inches(0.4),
                    k, font=SANS, size=12, bold=True, color=INK)
        add_textbox(slide, Inches(7.5), y + Inches(0.4), Inches(5.2), Inches(0.5),
                    v, font=SANS, size=12, color=SOFT, line_spacing=1.4)

    # Низ
    add_rect(slide, Inches(0.72), Inches(4.9), Inches(11.9), Inches(1.5),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(5.05), Inches(11.5), Inches(0.4),
                "ИТОГ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(5.4), Inches(11.5), Inches(0.9),
                ["Журналы бракеража, замен, выбора — как заполнены до 31.08.2026.",
                 "В сентябре — перепривязать документы к новому СанПиН, без паники."],
                font=SANS, size=14, color=INK, line_spacing=1.5)

    add_textbox(slide, Inches(0.72), Inches(6.55), Inches(12.0), Inches(0.4),
                "Источник S003  ·  пост. ГГСВ РФ от 02.06.2026 № 18, зарег. 02.06.2026 № 86854",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Доклад, вывод 1, раздел 17.12", 17)


# ============== СЛАЙД 18: ПИЛОТ — ФИНАЛ ==============
def slide_18_pilot_final(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)  # тёмный фон — единственный в основном маршруте

    # Терракотовая черта сверху
    add_rect(slide, Inches(0.72), Inches(0.85), Inches(1.0), Inches(0.06), fill=ACCENT)

    # Заголовок
    add_textbox(slide, Inches(0.72), Inches(1.1), Inches(12.0), Inches(0.6),
                "Пилот", font=SANS, size=12, bold=True, color=ACCENT)

    # Огромная цифра
    add_textbox(slide, Inches(0.72), Inches(1.7), Inches(8.0), Inches(2.0),
                "90\u00a0дней",
                font=SERIF, size=80, bold=True, color=ACCENT, line_spacing=1.0)

    # Под числом
    add_textbox(slide, Inches(0.72), Inches(3.7), Inches(7.5), Inches(0.5),
                "одно отделение  ·  одна неделя замера",
                font=SERIF, size=20, italic=True, color=IVORY, line_spacing=1.3)

    # Правая колонка — три строки
    add_textbox(slide, Inches(8.5), Inches(1.7), Inches(4.2), Inches(0.4),
                "ЧТО ДЕЛАЕМ", font=SANS, size=10, bold=True, color=ACCENT)

    items = [
        ("Два варианта", "на завтрак и обед"),
        ("Стандартный", "гарантирован всегда"),
        ("Житель выбрал —", "житель получил"),
    ]
    for i, (k, v) in enumerate(items):
        y = Inches(2.1 + i * 0.9)
        add_textbox(slide, Inches(8.5), y, Inches(4.2), Inches(0.5),
                    k, font=SERIF, size=18, bold=True, color=IVORY, line_spacing=1.2)
        add_textbox(slide, Inches(8.5), y + Inches(0.45), Inches(4.2), Inches(0.4),
                    v, font=SANS, size=13, color=GRAY, line_spacing=1.3)

    # Финальная строка
    add_accent_bar(slide, Inches(0.72), Inches(5.5), w=Inches(1.0), h=Inches(0.04))

    add_textbox(slide, Inches(0.72), Inches(5.7), Inches(12.0), Inches(0.6),
                "Измерить  →  проверить  →  скорректировать  →  решить о масштабе",
                font=SERIF, size=22, bold=True, color=IVORY, line_spacing=1.3)

    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.5),
                "Не внедряйте вслепую.  Измеряйте.",
                font=SERIF, size=18, italic=True, color=ACCENT, line_spacing=1.3)

    # Footer
    add_textbox(slide, Inches(0.72), Inches(7.1), Inches(10.0), Inches(0.25),
                "Доклад, раздел 35 (пилот), 22 (панель чисел)",
                font=SANS, size=9, color=GRAY)
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(1.1), Inches(0.25),
                "18 / 18", font=SANS, size=9, color=GRAY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(0.3), Inches(0.25),
                "18", font=SANS, size=9, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


# ============== СЛАЙД 19: РАЗДЕЛИТЕЛЬ — РЕЗЕРВ ==============
def slide_19_separator(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Резерв · не показывать в основном ходе")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Основной показ закончен.",
                font=SERIF, size=36, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.7), Inches(12.0), Inches(0.5),
                "Дальше — резерв для ответов.",
                font=SERIF, size=20, italic=True, color=SOFT)

    add_rect(slide, Inches(0.72), Inches(2.7), Inches(11.9), Inches(3.8),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.9), Inches(11.5), Inches(0.4),
                "ЧТО ЗДЕСЬ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(3.3), Inches(11.5), Inches(3.2),
                ["R1  разогрев  ·  вчерашний ужин, кто мог выбрать",
                 "R2  три конфликта программы",
                 "R3  Успенский ПНИ  ·  М2 — два вторых ежедневно",
                 "R4  Болотнинский ПНИ  ·  М1 — выбор несколько раз в неделю",
                 "R5  Усть-Илимский ДСО",
                 "R6  подробный кейс «Серафимовского»",
                 "R7  уровни IDDSI  ·  язык описания текстур",
                 "R8  правовая цепочка",
                 "R9  пищевая карта  ·  R10  журнал выбора",
                 "R11  закупочная модель  ·  R12  пять измеряемых показателей"],
                font=SANS, size=11, color=INK, line_spacing=1.4)

    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "На 90 минут достаточно слайдов 1–18.  Остальное — если зал спросит.",
                font=SERIF, size=15, italic=True, color=SOFT, align=PP_ALIGN.CENTER)

    add_footer(slide, "Приложение после финального слайда", "—", total="—")


# ============== СБОРКА ==============
def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 18 основных
    slide_01_title(prs)
    slide_02_five_dimensions(prs)
    slide_03_dysphagia(prs)
    slide_04_2190(prs)
    slide_05_74(prs)
    slide_06_night(prs)
    slide_07_folder(prs)
    slide_08_irkutsk(prs)
    slide_09_serafimovsky(prs)
    slide_10_order(prs)
    slide_11_procurement(prs)
    slide_12_doctor_first(prs)
    slide_13_model(prs)
    slide_14_channels(prs)
    slide_15_five_steps(prs)
    slide_16_three_errors(prs)
    slide_17_sanpin(prs)
    slide_18_pilot_final(prs)

    # Разделитель
    slide_19_separator(prs)

    return prs


if __name__ == "__main__":
    HERE = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(HERE, "Презентация_v6_Не_просто_накормить.pptx")
    prs = build_presentation()

    # Добавляем резервные слайды R1-R12
    from reserve_v6 import add_reserve_slides
    add_reserve_slides(prs)

    prs.save(out_path)
    print(f"PPTX saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
