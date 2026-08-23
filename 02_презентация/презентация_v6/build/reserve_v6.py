# -*- coding: utf-8 -*-
"""
reserve_v6.py — 12 резервных слайдов R1-R12 для v6.
Подключается к build_v6.py через функцию add_reserve_slides(prs).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ============== ПАЛИТРА (копия из build_v6.py) ==============
IVORY      = RGBColor(0xFA, 0xF9, 0xF5)
PANEL      = RGBColor(0xF0, 0xEE, 0xE6)
PANELWARM  = RGBColor(0xF6, 0xF1, 0xE9)
ACCENTSOFT = RGBColor(0xF5, 0xE7, 0xDE)
INK        = RGBColor(0x14, 0x14, 0x13)
BODY       = RGBColor(0x3D, 0x39, 0x29)
MUTED      = RGBColor(0x87, 0x86, 0x7F)
SOFT       = RGBColor(0x6E, 0x6A, 0x5E)
ACCENT     = RGBColor(0xC1, 0x5F, 0x3C)
ACCENTDEEP = RGBColor(0xB8, 0x55, 0x2F)
LINE       = RGBColor(0xE3, 0xDC, 0xCE)
GRAY       = RGBColor(0xB7, 0xAF, 0xA0)

SERIF = "Playfair Display"
SANS  = "PT Sans"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============== HELPERS ==============
def add_textbox(slide, x, y, w, h, text, *,
                font=SANS, size=14, bold=False, italic=False,
                color=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.3, word_wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    if isinstance(text, str): text = [text]
    for i, line in enumerate(text):
        if i == 0: p = tf.paragraphs[0]
        else: p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=0, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        try: shape.adjustments[0] = corner_radius / 100.0
        except Exception: pass
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width) if line_width else Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, weight_pt=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_kicker(slide, label):
    add_textbox(slide, Inches(0.72), Inches(0.5), Inches(8.0), Inches(0.3),
                label.upper(), font=SANS, size=11, bold=True, color=ACCENT)
    add_line(slide, Inches(0.72) + Inches(2.6), Inches(0.65), Inches(12.6), Inches(0.65),
             color=LINE, weight_pt=0.75)


def add_footer(slide, left_text, page_num):
    """Подвал резерва: только метка слайда (R1–R12)."""
    add_textbox(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.25),
                str(page_num), font=SANS, size=9, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


def set_bg(slide, color=IVORY):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def add_accent_bar(slide, x, y, w=Inches(0.8), h=Inches(0.05), color=ACCENT):
    return add_rect(slide, x, y, w, h, fill=color)


# ============== R1: РАЗОГРЕВ ==============
def slide_R1_warmup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Разогрев · 1 минута")

    # Кружок с ?
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(1.6),
                                     Inches(2.5), Inches(2.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = IVORY
    circle.line.color.rgb = LINE; circle.line.width = Pt(1.5)
    add_textbox(slide, Inches(0.95), Inches(1.6), Inches(2.5), Inches(2.5),
                "?", font=SERIF, size=140, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Линия
    add_rect(slide, Inches(4.0), Inches(1.95), Inches(0.06), Inches(2.0), fill=ACCENT)

    # Текст
    add_textbox(slide, Inches(4.3), Inches(1.85), Inches(8.5), Inches(2.0),
                "Вспомните вчерашний ужин.",
                font=SERIF, size=28, bold=True, color=INK, line_spacing=1.2)
    add_textbox(slide, Inches(4.3), Inches(2.55), Inches(8.5), Inches(2.0),
                "Поднимите руку, если хоть один житель вашего дома мог выбрать второе блюдо.",
                font=SERIF, size=18, italic=True, color=SOFT, line_spacing=1.3)

    add_textbox(slide, Inches(0.72), Inches(5.0), Inches(12.0), Inches(0.5),
                "Рука, которая не поднимается, — это и есть повестка наших 90 минут.",
                font=SERIF, size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.72), Inches(5.7), Inches(12.0), Inches(0.4),
                "Честное «никто не мог» — нормальный старт, а не провал.",
                font=SANS, size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, "Вопрос залу", "R1")


# ============== R2: ТРИ КОНФЛИКТА ПРОГРАММЫ ==============
def slide_R2_three_conflicts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Три конфликта программы")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Шесть вопросов программы — три конфликта.",
                font=SERIF, size=28, bold=True, color=INK)

    cards = [
        ("①", "Выбор  ↔  безопасность", "вопросы 1 и 5",
         "Субъектность жителя, недееспособность, медицинские рамки."),
        ("②", "Вариативность  ↔  норматив и бюджет", "вопросы 2, 4 и 6",
         "Меню в нормах, „лечебность“ руками кухни, качество против цены."),
        ("③", "Единое меню  ↔  индивидуальная поддержка", "вопрос 3",
         "Повар назначает диету, или медицинская служба?"),
    ]
    for i, (num, head, qs, body) in enumerate(cards):
        y = Inches(2.0 + i * 1.5)
        add_rect(slide, Inches(0.72), y, Inches(11.9), Inches(1.3),
                 fill=PANEL, corner_radius=12)
        add_textbox(slide, Inches(0.95), y + Inches(0.2), Inches(0.6), Inches(0.7),
                    num, font=SERIF, size=36, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.7), y + Inches(0.2), Inches(7.5), Inches(0.5),
                    head, font=SERIF, size=18, bold=True, color=INK)
        add_textbox(slide, Inches(1.7), y + Inches(0.7), Inches(7.5), Inches(0.5),
                    body, font=SANS, size=12, color=SOFT)
        add_textbox(slide, Inches(9.3), y + Inches(0.4), Inches(3.0), Inches(0.4),
                    qs, font=SANS, size=11, italic=True, color=ACCENTDEEP,
                    align=PP_ALIGN.RIGHT)

    add_footer(slide, "Программа организаторов; речь 0:00–2:00", "R2")


# ============== R3: УСПЕНСКИЙ ПНИ (М2) ==============
def slide_R3_uspenskiy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Модель М2 · Успенский ПНИ")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Успенский ПНИ · Новосибирская область",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Выбор по всем позициям ежедневно, во всех отделениях.",
                font=SERIF, size=18, italic=True, color=SOFT)

    add_rect(slide, Inches(0.72), Inches(2.5), Inches(11.9), Inches(2.8),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.7), Inches(11.5), Inches(0.4),
                "КОНФИГУРАЦИЯ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(3.1), Inches(11.5), Inches(2.0),
                ["Ежедневно:",
                 "  ·  два первых блюда",
                 "  ·  два вторых блюда",
                 "  ·  два салата",
                 "  ·  два напитка"],
                font=SERIF, size=16, color=INK, line_spacing=1.3)

    add_textbox(slide, Inches(0.72), Inches(5.5), Inches(12.0), Inches(0.5),
                "Диетические ограничения — внутри того же выбора, по назначению врача.",
                font=SANS, size=14, italic=True, color=SOFT)
    add_textbox(slide, Inches(0.72), Inches(5.7), Inches(12.0), Inches(0.4),
                "Не путать с ДДСОЛ (там — буфет, S019).",
                font=SANS, size=12, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "Источник S016  ·  upni.nso.ru/news/1916  ·  13.08.2024",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Доклад, раздел 25, лестница М0–М5", "R3")


# ============== R4: БОЛОТНИНСКИЙ ПНИ (М1) ==============
def slide_R4_bolotninskiy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Модель М1 · Болотнинский ПНИ")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Болотнинский ПНИ · Новосибирская область",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Выбор на части приёмов, несколько раз в неделю.",
                font=SERIF, size=18, italic=True, color=SOFT)

    add_rect(slide, Inches(0.72), Inches(2.5), Inches(11.9), Inches(2.5),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.7), Inches(11.5), Inches(0.4),
                "КОНФИГУРАЦИЯ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(3.1), Inches(11.5), Inches(1.8),
                ["Завтрак и обед — выбор ключевых позиций.",
                 "",
                 "Перечень определяет диетсестра.",
                 "Подходит как стартовая модель."],
                font=SERIF, size=18, color=INK, line_spacing=1.4)

    add_textbox(slide, Inches(0.72), Inches(5.3), Inches(12.0), Inches(0.5),
                "Не путать с Успенским ПНИ. Это не «Новосибирский ПНИ» как отдельный дом.",
                font=SANS, size=12, italic=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "Источник S015  ·  bpni.nso.ru/news/2312  ·  16.08.2024",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Доклад, раздел 25, лестница М0–М5", "R4")


# ============== R5: УСТЬ-ИЛИМСК ==============
def slide_R5_ustilimsk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Модель · Усть-Илимский ДСО")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Усть-Илимский ДСО · Иркутская область",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Публикация 27.09.2024: «Что за зверь такой? Вариативное питание».",
                font=SERIF, size=18, italic=True, color=SOFT)

    add_rect(slide, Inches(0.72), Inches(2.5), Inches(11.9), Inches(2.5),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.7), Inches(11.5), Inches(0.4),
                "ЧТО ОПИСАНО", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(3.1), Inches(11.5), Inches(1.8),
                ["Дом ведёт публичный рассказ о внедрении",
                 "вариативного питания — для коллег и жителей.",
                 "",
                 "Один из четырёх публичных кейсов в России."],
                font=SERIF, size=16, color=INK, line_spacing=1.4)

    add_textbox(slide, Inches(0.72), Inches(5.3), Inches(12.0), Inches(0.5),
                "Публикация — не аудит эффекта. Источник первичный, открытый.",
                font=SANS, size=12, italic=True, color=SOFT)
    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "Источник S017  ·  ui-ogbuso.ru  ·  27.09.2024",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Доклад, раздел 25", "R5")


# ============== R6: ПОДРОБНЫЙ КЕЙС СЕРАФИМОВСКОГО ==============
def slide_R6_serafimovsky_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Подробно · ДСО «Серафимовский»")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Жители приняли выбор с энтузиазмом.",
                font=SERIF, size=26, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Несколько месяцев часть жителей выбирает рацион накануне.",
                font=SERIF, size=18, italic=True, color=SOFT)

    add_rect(slide, Inches(0.72), Inches(2.5), Inches(11.9), Inches(1.5),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.65), Inches(11.5), Inches(0.4),
                "ФОРМАТ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(3.0), Inches(11.5), Inches(0.9),
                ["«Ресторанные дни» по отделениям — рассказ зам. директора И. Ю. Игнатьевой,",
                 "посещение Уполномоченного по правам человека в СПб 19.05.2026."],
                font=SANS, size=13, color=INK, line_spacing=1.45)

    add_rect(slide, Inches(0.72), Inches(4.2), Inches(11.9), Inches(1.4),
             fill=ACCENTSOFT, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(4.35), Inches(11.5), Inches(0.4),
                "ЧЕСТНО О ЦЕНЕ", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(4.7), Inches(11.5), Inches(0.9),
                ["Старт не был бесплатным — оборудование и ставки пищеблока.",
                 "Это сказано прямо: «рассказываю и про цену, и про пределы»."],
                font=SANS, size=13, color=INK, line_spacing=1.45)

    add_textbox(slide, Inches(0.72), Inches(5.85), Inches(12.0), Inches(0.4),
                "Старт — в автономном учреждении, 223-ФЗ, тендеры на tenmon.ru и tenderguru.ru.",
                font=SANS, size=12, italic=True, color=SOFT)
    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "Источник S030  ·  gov.spb.ru, 22.07.2025  ·  upchspb.ru/news/32938  ·  19.05.2026",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Доклад, раздел 25, 29", "R6")


# ============== R7: IDDSI ==============
def slide_R7_iddsi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Уровни IDDSI")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "IDDSI — язык описания текстур.",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Не диагностический протокол. Не замена назначения врача.",
                font=SERIF, size=18, italic=True, color=SOFT)

    levels = [
        ("0", "Жидкая", "тонкая"),
        ("1", "Слегка густая", "для части жидкостей"),
        ("2", "Слабо густая", ""),
        ("3", "Умеренно густая", "жидкое пюре"),
        ("4", "Очень густая", "пюре"),
        ("5", "Мягкая", "мини-кусочки, мягкая"),
        ("6", "Мягкая и размерами", "мягкие кусочки ≤ 1,5 см"),
        ("7", "Обычная", "легко жуётся"),
    ]
    bar_x = 0.72
    bar_w = 12.0
    n = len(levels)
    col_w = (bar_w - 0.5 * (n - 1)) / n

    for i, (num, name, desc) in enumerate(levels):
        x = Inches(bar_x + i * (col_w + 0.1))
        y = Inches(2.6)
        # Столбик
        h_in = 1.5 + (i * 0.15)  # визуальный рост
        add_rect(slide, x, y, Inches(col_w), Inches(1.0), fill=PANEL, corner_radius=8)
        # Номер
        add_textbox(slide, x, y + Inches(0.1), Inches(col_w), Inches(0.4),
                    f"{num}", font=SERIF, size=22, bold=True, color=ACCENT,
                    align=PP_ALIGN.CENTER)
        # Название
        add_textbox(slide, x, y + Inches(0.5), Inches(col_w), Inches(0.4),
                    name, font=SANS, size=10, bold=True, color=INK,
                    align=PP_ALIGN.CENTER)
        # Описание
        add_textbox(slide, x - Inches(0.05), y + Inches(1.05), Inches(col_w + 0.1), Inches(0.6),
                    desc, font=SANS, size=8, color=SOFT,
                    align=PP_ALIGN.CENTER, line_spacing=1.2)

    add_rect(slide, Inches(0.72), Inches(4.7), Inches(11.9), Inches(1.4),
             fill=ACCENTSOFT, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(4.85), Inches(11.5), Inches(0.4),
                "ОГОВОРКА", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(5.2), Inches(11.5), Inches(0.8),
                ["IDDSI — это терминология и методика проверки консистенции.",
                 "Текстуру назначает врач. Кухня исполняет по таблице."],
                font=SANS, size=13, color=INK, line_spacing=1.45)

    add_textbox(slide, Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.4),
                "International Dysphagia Diet Standardisation Initiative  ·  IDDSI Framework 2019",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Доклад, раздел 10  ·  IDDSI — не протокол лечения", "R7")


# ============== R8: ПРАВОВАЯ ЦЕПОЧКА ==============
def slide_R8_legal_chain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Правовая цепочка")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Четыре уровня правовой рамки.",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "От статуса жителя — до санитарного акта.",
                font=SERIF, size=18, italic=True, color=SOFT)

    items = [
        ("①", "3185-1 ст. 5 ч. 2; ст. 37, 43", "Закон РФ о психиатрической помощи",
         "Гуманное отношение — ст. 5 ч. 2, не ст. 37 ч. 2. Ст. 37 и 43 — права пациентов и статус проживающих."),
        ("②", "442-ФЗ ст. 9, 16, 32 ч. 4", "Основы социального обслуживания",
         "Уважительное и гуманное отношение; ИППСУ исходя из потребности; предельная величина платы 75 %."),
        ("③", "ГК РФ ст. 29, 30", "Гражданский кодекс",
         "Опекун недееспособного; учёт мнения подопечного."),
        ("④", "СанПиН 2.3/2.4.4282-26, п. 56(11)", "с 01.09.2026",
         "Не менее 3-х приёмов в день, в т. ч. диетическое (лечебное) по медицинским показаниям."),
    ]
    for i, (num, law, head, body) in enumerate(items):
        y = Inches(2.4 + i * 1.0)
        add_rect(slide, Inches(0.72), y, Inches(11.9), Inches(0.85),
                 fill=PANEL, corner_radius=10)
        add_textbox(slide, Inches(0.95), y + Inches(0.1), Inches(0.6), Inches(0.6),
                    num, font=SERIF, size=22, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.7), y + Inches(0.1), Inches(4.5), Inches(0.4),
                    law, font=SERIF, size=14, bold=True, color=INK)
        add_textbox(slide, Inches(1.7), y + Inches(0.45), Inches(4.5), Inches(0.4),
                    head, font=SANS, size=10, italic=True, color=SOFT)
        add_textbox(slide, Inches(6.4), y + Inches(0.15), Inches(6.0), Inches(0.6),
                    body, font=SANS, size=11, color=INK, line_spacing=1.35)

    add_footer(slide, "Доклад, раздел 17, 17.12", "R8")


# ============== R9: ПИЩЕВАЯ КАРТА ==============
def slide_R9_food_card(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Пример · пищевая карта жителя")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Пищевая карта — пять строк.",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Один житель. Один документ. Хватает для пищевой безопасности и пищевого выбора.",
                font=SERIF, size=16, italic=True, color=SOFT)

    rows = [
        ("Профиль", "диагноз, ИППСУ, ограничения, лекарства"),
        ("Текстура", "назначение врача, IDDSI-уровень"),
        ("Исключения", "аллергии, лекарственно-пищевые взаимодействия (клозапин ↔ кофеин)"),
        ("Предпочтения", "биография, что ел всю жизнь, к чему тянется"),
        ("Способ фиксации выбора", "показ, наблюдение, слово, сведения от знающих"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.5 + i * 0.7)
        add_rect(slide, Inches(0.72), y, Inches(11.9), Inches(0.6),
                 fill=PANEL, corner_radius=8)
        add_textbox(slide, Inches(0.95), y + Inches(0.1), Inches(3.5), Inches(0.4),
                    k, font=SERIF, size=14, bold=True, color=ACCENTDEEP)
        add_textbox(slide, Inches(4.5), y + Inches(0.1), Inches(8.0), Inches(0.4),
                    v, font=SANS, size=12, color=INK)

    add_textbox(slide, Inches(0.72), Inches(6.3), Inches(12.0), Inches(0.4),
                "Подпись врача на карте обязательна.",
                font=SERIF, size=12, italic=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.72), Inches(6.7), Inches(12.0), Inches(0.3),
                "Источник: доклад, приложение 3; пример из немецкого стандарта (§ 6 Abs. 4 SGB XII)",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Шаблон — проект, не норматив", "R9")


# ============== R10: ЖУРНАЛ ВЫБОРА ==============
def slide_R10_choice_journal(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Пример · журнал выбора")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Журнал выбора — одна строка на жителя в день.",
                font=SERIF, size=26, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Без журнала «выбор» — это только слова. Проверяющий это знает.",
                font=SERIF, size=16, italic=True, color=SOFT)

    headers = ["Дата", "Отделение", "Житель", "Заказал", "Получил", "Подпись смены"]
    rows = [
        ["26.08.2026", "№ 2", "Иванов И. И.", "Гречка", "Гречка", "Петрова"],
        ["26.08.2026", "№ 2", "Сидорова А. П.", "Рис", "Рис", "Петрова"],
        ["26.08.2026", "№ 2", "Ким В. С.", "—", "Стандарт", "Петрова"],
    ]

    table_x = 0.72
    table_y = 2.5
    col_widths = [1.5, 1.3, 2.4, 2.0, 2.0, 2.7]
    row_h = 0.5

    # Заголовки
    x = table_x
    for i, h in enumerate(headers):
        add_rect(slide, Inches(x), Inches(table_y), Inches(col_widths[i]), Inches(row_h),
                 fill=ACCENTDEEP, corner_radius=4)
        add_textbox(slide, Inches(x + 0.1), Inches(table_y + 0.05), Inches(col_widths[i] - 0.1), Inches(row_h - 0.1),
                    h, font=SANS, size=10, bold=True, color=IVORY, align=PP_ALIGN.CENTER)
        x += col_widths[i]

    # Строки
    for r, row in enumerate(rows):
        y = table_y + row_h + 0.05 + r * row_h
        x = table_x
        fill = PANEL if r % 2 == 0 else IVORY
        for i, cell in enumerate(row):
            add_rect(slide, Inches(x), Inches(y), Inches(col_widths[i]), Inches(row_h),
                     fill=fill, line_color=LINE, line_width=0.5, corner_radius=4)
            add_textbox(slide, Inches(x + 0.1), Inches(y + 0.1), Inches(col_widths[i] - 0.1), Inches(row_h - 0.2),
                        cell, font=SANS, size=10, color=INK, align=PP_ALIGN.CENTER)
            x += col_widths[i]

    # Подпись
    add_textbox(slide, Inches(0.72), Inches(5.4), Inches(12.0), Inches(0.5),
                "Не выбрал — стандартный вариант. Без давления.",
                font=SERIF, size=14, italic=True, color=ACCENTDEEP, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.72), Inches(6.0), Inches(12.0), Inches(0.4),
                "Источник: доклад, приложение 4",
                font=SANS, size=10, color=MUTED)

    add_footer(slide, "Тетрадный формат допустим на старте", "R10")


# ============== R11: ЗАКУПОЧНАЯ МОДЕЛЬ ==============
def slide_R11_procurement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Закупочная модель")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "223-ФЗ и 44-ФЗ: что важно знать директору.",
                font=SERIF, size=26, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Вариативность возможна в обоих контурах — но документы разные.",
                font=SERIF, size=18, italic=True, color=SOFT)

    # 223-ФЗ
    add_rect(slide, Inches(0.72), Inches(2.5), Inches(5.8), Inches(3.5),
             fill=PANEL, corner_radius=14)
    add_textbox(slide, Inches(0.95), Inches(2.7), Inches(5.5), Inches(0.4),
                "223-ФЗ", font=SERIF, size=22, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.95), Inches(3.2), Inches(5.5), Inches(0.4),
                "автономное учреждение, конкурсная основа",
                font=SANS, size=11, italic=True, color=SOFT)
    add_textbox(slide, Inches(0.95), Inches(3.7), Inches(5.5), Inches(2.0),
                ["·  Положение о закупке",
                 "·  Конкурсная документация",
                 "·  Тендеры на tenmon.ru, tenderguru.ru",
                 "·  Источник S030 для Серафимовского"],
                font=SANS, size=12, color=INK, line_spacing=1.5)

    # 44-ФЗ
    add_rect(slide, Inches(6.8), Inches(2.5), Inches(5.8), Inches(3.5),
             fill=ACCENTSOFT, corner_radius=14)
    add_textbox(slide, Inches(7.05), Inches(2.7), Inches(5.5), Inches(0.4),
                "44-ФЗ", font=SERIF, size=22, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(7.05), Inches(3.2), Inches(5.5), Inches(0.4),
                "бюджетное учреждение, госзакупка",
                font=SANS, size=11, italic=True, color=SOFT)
    add_textbox(slide, Inches(7.05), Inches(3.7), Inches(5.5), Inches(2.0),
                ["·  Контрактная система",
                 "·  zakupki.gov.ru",
                 "·  Источник для Тесового берега:",
                 "    ИНН 7827661874, ЕГРЮЛ"],
                font=SANS, size=12, color=INK, line_spacing=1.5)

    add_textbox(slide, Inches(0.72), Inches(6.3), Inches(12.0), Inches(0.5),
                "Что в обоих случаях: положение о питании, журнал заказа, согласование с учредителем.",
                font=SERIF, size=14, italic=True, color=SOFT, align=PP_ALIGN.CENTER)

    add_footer(slide, "Доклад, раздел 29", "R11")


# ============== R12: ПЯТЬ ИЗМЕРЯЕМЫХ ПОКАЗАТЕЛЕЙ ==============
def slide_R12_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)
    add_kicker(slide, "Пять измеряемых показателей")

    add_textbox(slide, Inches(0.72), Inches(0.95), Inches(12.0), Inches(0.7),
                "Пилот измеряет пять вещей.",
                font=SERIF, size=28, bold=True, color=INK)
    add_textbox(slide, Inches(0.72), Inches(1.5), Inches(12.0), Inches(0.5),
                "Без замера — нет ответа на вопрос «стало лучше?».",
                font=SERIF, size=18, italic=True, color=SOFT)

    metrics = [
        ("①", "Пищевая карта", "заполнена у каждого жителя"),
        ("②", "Доля с реальным выбором", "приёмы, где альтернатива была доступна"),
        ("③", "Исполненные заказы", "житель выбрал — житель получил"),
        ("④", "Отказы, замены, отходы", "что не съели, что заменили, что ушло в отходы"),
        ("⑤", "События риска", "поперхивания, пропуски, потеря веса, обострения"),
    ]
    for i, (num, head, body) in enumerate(metrics):
        y = Inches(2.4 + i * 0.65)
        add_rect(slide, Inches(0.72), y, Inches(11.9), Inches(0.55),
                 fill=PANEL, corner_radius=8)
        add_textbox(slide, Inches(0.95), y + Inches(0.1), Inches(0.6), Inches(0.4),
                    num, font=SERIF, size=20, bold=True, color=ACCENT)
        add_textbox(slide, Inches(1.7), y + Inches(0.1), Inches(4.0), Inches(0.4),
                    head, font=SERIF, size=14, bold=True, color=INK)
        add_textbox(slide, Inches(5.8), y + Inches(0.1), Inches(6.5), Inches(0.4),
                    body, font=SANS, size=12, color=SOFT)

    add_rect(slide, Inches(0.72), Inches(5.85), Inches(11.9), Inches(0.7),
             fill=ACCENTSOFT, corner_radius=10)
    add_textbox(slide, Inches(0.95), Inches(5.95), Inches(11.5), Inches(0.4),
                "ВАЖНО", font=SANS, size=10, bold=True, color=ACCENTDEEP)
    add_textbox(slide, Inches(0.95), Inches(6.2), Inches(11.5), Inches(0.4),
                "Процент выбравших — не KPI.  Человек вправе выбрать стандартное или не выбирать вовсе.",
                font=SERIF, size=14, bold=True, italic=True, color=INK)

    add_footer(slide, "Доклад, раздел 22 (панель чисел)", "R12")


# ============== СБОРКА ==============
def add_reserve_slides(prs):
    """Добавить 12 резервных слайдов к уже собранной презентации."""
    slide_R1_warmup(prs)
    slide_R2_three_conflicts(prs)
    slide_R3_uspenskiy(prs)
    slide_R4_bolotninskiy(prs)
    slide_R5_ustilimsk(prs)
    slide_R6_serafimovsky_detail(prs)
    slide_R7_iddsi(prs)
    slide_R8_legal_chain(prs)
    slide_R9_food_card(prs)
    slide_R10_choice_journal(prs)
    slide_R11_procurement(prs)
    slide_R12_metrics(prs)
